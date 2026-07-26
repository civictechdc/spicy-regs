"""Real-provider tests for the Docling adapter, against the exact pinned releases.

These run only where the optional extra is installed, and they never invoke model
inference: they build real converters and real ``docling-core`` documents, and
convert only Office files, which ``SimplePipeline`` handles with a declarative
backend and no model at all. Run them with::

    uv run --frozen --extra docling pytest tests/test_docpipeline_adapter_docling_real.py

What they exist to hold:

* the release blocker — a table-only DOCX has one ``TableItem`` with no ``.text``
  at all, so an adapter that reads ``.text`` completes with every cell value
  dropped;
* all three supported formats end to end: DOCX for the table blocker, PPTX for
  speaker notes, and XLSX through the real Excel backend;
* the exact ``docling-core`` content-layer enum, that PowerPoint speaker notes
  really arrive on the ``notes`` layer the default iteration hides, and that
  ``background`` and ``invisible`` text survives too — with each layer's heading
  context staying inside that layer;
* every closed provider-token set the adapter validates against — item labels,
  conversion statuses, input formats, coordinate origins, and the table-bearing
  labels — held *complete and exactly spelled* against the pinned enums
  themselves, so a remembered spelling or a missing member cannot drift from the
  release;
* the exact ``TableCell`` geometry, including a real merged cell, real cell text
  holding the serialization's own separators, real declared shapes that
  contradict their cells, a real grid empty in one dimension only, a real cell
  whose span attribute is gone, and both directions of the label/shape
  agreement — a table-labelled item with no readable ``TableData``, and real
  ``TableData`` carried under a label that is not a table;
* the ``docling-core`` formula and caption behavior the adapter depends on: that
  ``add_table`` accepts any ``TextItem`` as a caption, that two tables may share
  one caption target, and that a dangling, self-naming, or twice-named reference
  still fails closed; and
* that the loaded converter really builds the Office-only, model-free path, that
  the policy records every bound it enforces — at the scope each name states —
  and the containment it does not, and that a PDF is refused by name before any
  of it runs.
"""

from __future__ import annotations

import hashlib
import importlib
import inspect
import io
import pkgutil
import re
import typing
import warnings
from pathlib import Path
from types import SimpleNamespace
from typing import Any

import pytest

docling = pytest.importorskip("docling", reason="needs the optional docling extra")
docling_core = pytest.importorskip("docling_core", reason="needs the optional docling extra")

from docling.datamodel.base_models import (  # noqa: E402  # ty: ignore[unresolved-import]
    ConversionStatus,
    InputFormat,
)
from docling.datamodel.pipeline_options import (  # noqa: E402  # ty: ignore[unresolved-import]
    ConvertPipelineOptions,
)
from docling.document_converter import (  # noqa: E402  # ty: ignore[unresolved-import]
    DocumentConverter,
    ExcelFormatOption,
    PowerpointFormatOption,
    WordFormatOption,
)
from docling.exceptions import ConversionError  # noqa: E402  # ty: ignore[unresolved-import]
from docling.pipeline.simple_pipeline import SimplePipeline  # noqa: E402  # ty: ignore[unresolved-import]
from docling_core.types.doc import (  # noqa: E402  # ty: ignore[unresolved-import]
    BoundingBox,
    CoordOrigin,
    ProvenanceItem,
)
from docling_core.types.doc.common.content_layer import (  # noqa: E402  # ty: ignore[unresolved-import]
    ContentLayer,
)
from docling_core.types.doc import document as document_module  # noqa: E402  # ty: ignore[unresolved-import]
from docling_core.types.doc.document import (  # noqa: E402  # ty: ignore[unresolved-import]
    DoclingDocument,
    RefItem,
    SectionHeaderItem,
    TableCell,
    TableData,
    TableItem,
)
from docling_core.types.doc.labels import DocItemLabel  # noqa: E402  # ty: ignore[unresolved-import]

from spicy_regs.docpipeline.adapters import docling as adapter_module  # noqa: E402
from spicy_regs.docpipeline.adapters.docling import (  # noqa: E402
    ACCEPTED_CONVERSION_STATUSES,
    ADAPTER_MAPPING_REVISION,
    CONTENT_FROM_ORIG,
    CONTENT_FROM_TABLE_CELLS,
    CONTENT_FROM_TEXT,
    CONTENT_LAYERS,
    CONVERSION_STATUSES,
    COORDINATE_ORIGINS,
    DOC_ITEM_LABELS,
    DOCLING_CORE_VERSION,
    DOCLING_VERSION,
    ELEMENT_SEPARATOR,
    FALLBACK_ERROR_TYPE,
    FORMAT_DOCX,
    FORMAT_PPTX,
    FORMAT_XLSX,
    MAX_CAPTION_REFS_PER_ITEM,
    MAX_CELLS_PER_TABLE,
    MAX_ERROR_TYPE_CHARS,
    MAX_HEADING_LEVEL,
    MAX_ITEMS,
    MAX_MAPPED_CHARACTERS,
    MAX_PAGE_NUMBER,
    MAX_PROVENANCE_CHARACTER_INDEX,
    MAX_PROVIDER_ERRORS,
    MAX_PROVIDER_TOKEN_CHARS,
    MAX_REFERENCE_CHARS,
    MAX_REGIONS_PER_ITEM,
    MAX_TABLE_CELL_CHARACTERS,
    MAX_TABLE_DIMENSION,
    MAX_TABLES,
    MAX_TOTAL_CAPTION_REFS,
    MAX_TOTAL_REGIONS,
    MAX_TOTAL_TABLE_CELLS,
    MAX_TREE_DEPTH,
    NO_CONTENT,
    NO_COORDINATES,
    PARSED_TEXT_OFFSETS,
    PARSER_PAGE_COORDINATES,
    PIPELINE_SIMPLE,
    PROVIDER_INPUT_FORMATS,
    SOURCE_NAME_ENCODING,
    SUPPORTED_FORMATS,
    TABLE_KINDS,
    TABLE_SERIALIZATION,
    DoclingDocumentParser,
    DoclingParseError,
    bounded_error_type,
    installed_package_version,
)

TABLE_ROWS = (("Pollutant", "Limit"), ("BOD5", "30 mg/L"))
TABLE_TEXT = "Pollutant\tLimit\nBOD5\t30 mg/L"
FORMULA_SOURCE = "C = m / V"
SPEAKER_NOTE = "Speaker note: confirm the 30 mg/L figure with OW."
SLIDE_TITLE = "Effluent Guidelines"


@pytest.fixture(scope="module", autouse=True)
def pinned_releases() -> None:
    """Refuse to interpret these results under any other release."""
    assert installed_package_version("docling") == DOCLING_VERSION
    assert installed_package_version("docling-core") == DOCLING_CORE_VERSION


def table_only_docx() -> bytes:
    """One 2x2 table and nothing else — the document that used to parse to nothing."""
    import docx  # ty: ignore[unresolved-import]

    document = docx.Document()
    table = document.add_table(rows=2, cols=2)
    for row, values in enumerate(TABLE_ROWS):
        for column, value in enumerate(values):
            table.cell(row, column).text = value
    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def merged_table_docx() -> bytes:
    """A merged header cell over two columns, plus cell text holding tabs and newlines."""
    import docx  # ty: ignore[unresolved-import]

    document = docx.Document()
    merged = document.add_table(rows=2, cols=2)
    merged.cell(0, 0).merge(merged.cell(0, 1)).text = "Spanning header"
    merged.cell(1, 0).text = "Line\twith\ttabs"
    merged.cell(1, 1).text = "two\nlines"
    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def headings_docx() -> bytes:
    """Two numbered headings, so the release really emits a ``SectionHeaderItem.level``."""
    import docx  # ty: ignore[unresolved-import]

    document = docx.Document()
    document.add_heading("Scope", level=1)
    document.add_paragraph("This part applies to discharges.")
    document.add_heading("Limits", level=2)
    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def limits_xlsx() -> bytes:
    """One sheet holding the same 2x2 grid, through the Excel backend."""
    from openpyxl import Workbook  # ty: ignore[unresolved-import]

    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Limits"
    for values in TABLE_ROWS:
        sheet.append(list(values))
    buffer = io.BytesIO()
    workbook.save(buffer)
    return buffer.getvalue()


def notes_pptx() -> bytes:
    """One slide with a title, a body box, and a speaker note."""
    from pptx import Presentation  # ty: ignore[unresolved-import]
    from pptx.util import Inches  # ty: ignore[unresolved-import]

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = SLIDE_TITLE
    box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1))
    box.text_frame.text = "Body bullet on the slide."
    slide.notes_slide.notes_text_frame.text = SPEAKER_NOTE
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


class RecordedConverter:
    """Wrap a real ``DocumentConverter`` to capture what the adapter asked of it."""

    def __init__(self, converter: Any) -> None:
        self._converter = converter
        self.calls: list[dict[str, Any]] = []

    def convert(self, source: Any, **options: Any) -> Any:
        self.calls.append({"name": Path(str(source)).name, "options": dict(options)})
        return self._converter.convert(source, **options)


def office_converter() -> Any:
    """A real converter for the Office formats, built without this adapter."""
    return DocumentConverter(
        allowed_formats=[InputFormat.DOCX, InputFormat.PPTX, InputFormat.XLSX],
        format_options={
            InputFormat.DOCX: WordFormatOption(),
            InputFormat.PPTX: PowerpointFormatOption(),
            InputFormat.XLSX: ExcelFormatOption(),
        },
    )


# --- the release blocker ----------------------------------------------------


def test_docling_core_keeps_table_values_only_in_data() -> None:
    conversion = DocumentConverter(allowed_formats=[InputFormat.DOCX]).convert(
        _stream(table_only_docx(), "table.docx"), raises_on_error=False
    )

    assert str(conversion.status.value) == "success"
    items = [item for item, _ in conversion.document.iterate_items()]
    assert [type(item).__name__ for item in items] == ["TableItem"]
    # No ``text`` attribute at all, and no page geometry for a declarative backend.
    assert not hasattr(items[0], "text")
    table = items[0]
    assert isinstance(table, TableItem)
    assert table.prov == []
    assert conversion.document.pages == {}
    assert [cell.text for cell in table.data.table_cells] == ["Pollutant", "Limit", "BOD5", "30 mg/L"]


def test_a_table_only_docx_parses_to_every_cell_value() -> None:
    result = DoclingDocumentParser().parse(table_only_docx(), source_name="effluent-limits.docx")

    assert result.document.text == "Pollutant\tLimit\nBOD5\t30 mg/L"
    assert [one.kind for one in result.document.elements] == ["table"]
    element = result.document.elements[0]
    assert (element.content_source, element.text_usable) == (CONTENT_FROM_TABLE_CELLS, True)
    assert (element.coordinate_grade, element.content_layer) == (NO_COORDINATES, "body")
    assert [
        (cell.row_start, cell.row_end, cell.column_start, cell.column_end, cell.text)
        for cell in result.document.tables[0].cells
    ] == [
        (0, 1, 0, 1, "Pollutant"),
        (0, 1, 1, 2, "Limit"),
        (1, 2, 0, 1, "BOD5"),
        (1, 2, 1, 2, "30 mg/L"),
    ]
    assert [cell.column_header for cell in result.document.tables[0].cells] == [True, True, False, False]
    assert result.document.tables[0].serialization_ambiguous is False

    call = result.call
    assert call.status == "completed"
    assert call.usable_character_count == len(result.document.text)
    assert (call.table_count, call.table_cell_count, call.omission_count) == (1, 4, 0)
    assert (call.input_format, call.provider_input_format) == (FORMAT_DOCX, FORMAT_DOCX)
    assert call.policy.pipeline == PIPELINE_SIMPLE
    assert call.policy.converter_source == "loaded"
    assert (call.coordinate_grade, call.content_layers_present) == (NO_COORDINATES, ("body",))

    # Offsets into the adapter's own text, exact and round-trippable.
    assert result.document.text[element.start_char : element.end_char] == element.text


def test_a_real_xlsx_parses_to_every_cell_with_a_complete_call_record() -> None:
    """The third supported format, end to end, through the real Excel backend.

    XLSX is the format with no fixture of its own before this: DOCX proved the
    table blocker and PPTX proved the notes layer, leaving the spreadsheet backend
    asserted only through the loaded converter's options.
    """
    result = DoclingDocumentParser().parse(limits_xlsx(), source_name="effluent-limits.xlsx")

    # The provider chose the Excel backend, and the record says so.
    assert (result.call.input_format, result.call.provider_input_format) == (FORMAT_XLSX, FORMAT_XLSX)
    assert result.document.input_format == FORMAT_XLSX

    # Every cell value survives as data, in reading order, plus the labeled
    # serialization — no value exists only in the rendering.
    assert len(result.document.tables) == 1
    table = result.document.tables[0]
    assert [cell.text for cell in table.cells] == ["Pollutant", "Limit", "BOD5", "30 mg/L"]
    assert (table.row_count, table.column_count) == (2, 2)
    assert [(cell.row_start, cell.row_end, cell.column_start, cell.column_end) for cell in table.cells] == [
        (0, 1, 0, 1),
        (0, 1, 1, 2),
        (1, 2, 0, 1),
        (1, 2, 1, 2),
    ]
    assert all(isinstance(cell.column_header, bool) and isinstance(cell.row_header, bool) for cell in table.cells)
    assert (table.serialization, table.serialization_ambiguous) == (TABLE_SERIALIZATION, False)
    assert table.content_layer == "body"
    assert TABLE_TEXT in result.document.text

    # Offsets address the adapter's own text, and every element round-trips.
    element = next(one for one in result.document.elements if one.parser_ref == table.parser_ref)
    assert (element.content_source, element.text_usable) == (CONTENT_FROM_TABLE_CELLS, True)
    assert element.content_layer == "body"
    assert result.document.text[element.start_char : element.end_char] == element.text == TABLE_TEXT
    for one in result.document.elements:
        assert result.document.text[one.start_char : one.end_char] == one.text
    assert result.document.offsets == PARSED_TEXT_OFFSETS

    call = result.call
    assert (call.status, call.failure_reason, call.error_type) == ("completed", None, None)
    assert (call.provider_invoked, call.attempt_count) == (True, 1)
    assert (call.conversion_status, call.provider_error_count) == ("success", 0)
    assert (call.table_count, call.table_cell_count) == (1, 4)
    assert call.content_layers_present == ("body",)
    assert call.usable_character_count == sum(len(one.text) for one in result.document.elements if one.text_usable)
    assert call.character_count == len(result.document.text)
    assert (call.policy.pipeline, call.policy.converter_source) == (PIPELINE_SIMPLE, "loaded")
    assert call.parser_id.endswith(call.policy_digest[:16])
    assert ADAPTER_MAPPING_REVISION in call.parser_id
    # Whatever geometry the Excel backend did or did not give, it is graded
    # consistently and never as source-exact evidence.
    assert call.coordinate_grade == (PARSER_PAGE_COORDINATES if element.regions else NO_COORDINATES)
    assert {
        region.coordinate_origin for one in result.document.elements for region in one.regions
    } <= COORDINATE_ORIGINS


# --- content layers ---------------------------------------------------------


def test_the_adapter_requests_every_content_layer_the_pinned_core_exposes() -> None:
    assert CONTENT_LAYERS == tuple(member.value for member in ContentLayer)
    # Docling's own default hides four of the five.
    from docling_core.types.doc.document import (  # ty: ignore[unresolved-import]
        DEFAULT_CONTENT_LAYERS,
    )

    assert {layer.value for layer in DEFAULT_CONTENT_LAYERS} == {"body"}


def test_the_closed_provider_token_sets_are_the_pinned_releases_own_members() -> None:
    """The pinned provider is the authority for every set the adapter closes.

    A hand-written constant that drifted from an enum would either refuse valid
    provider output or record a value the adapter does not actually understand, so
    every set is held to the enum itself — complete, and spelled exactly — rather
    than to a remembered convention. Each of these is the whole member list, not a
    sample: a missing member is a valid document this adapter would refuse.
    """
    assert COORDINATE_ORIGINS == {member.value for member in CoordOrigin}
    assert DOC_ITEM_LABELS == {member.value for member in DocItemLabel}
    assert CONVERSION_STATUSES == {member.value for member in ConversionStatus}
    assert PROVIDER_INPUT_FORMATS == {member.value for member in InputFormat}
    assert ACCEPTED_CONVERSION_STATUSES == {ConversionStatus.SUCCESS.value}
    assert ACCEPTED_CONVERSION_STATUSES < CONVERSION_STATUSES
    # The case is the enum's own, and it is not one convention across the release:
    # labels, statuses, and formats are lowercase, coordinate origins are not.
    # Lowercasing whatever arrives would have produced values from none of them.
    for closed in (DOC_ITEM_LABELS, CONVERSION_STATUSES, PROVIDER_INPUT_FORMATS):
        assert closed == {token.lower() for token in closed}
    assert COORDINATE_ORIGINS == {token.upper() for token in COORDINATE_ORIGINS}
    # Stated independently of the adapter's own pattern: every label the pinned
    # release can emit is a short lowercase token, so the one remaining lenient
    # reader (failure categories) cannot refuse valid provider output either.
    rejected = sorted(label for label in DOC_ITEM_LABELS if not re.fullmatch(r"[a-z0-9_-]{1,40}", label))
    assert not rejected
    assert MAX_PROVIDER_TOKEN_CHARS >= max(len(label) for label in DOC_ITEM_LABELS)


def test_the_table_bearing_kinds_are_exactly_the_labels_the_pinned_table_item_declares() -> None:
    """``TableItem`` is the only pinned item class that holds a grid, under two labels.

    Both directions of the adapter's label/shape agreement come from here: the
    ``Literal`` fixes which labels must carry ``TableData``, and the fact that no
    other item class declares a ``data`` field at all fixes which labels must not.
    """
    declared = typing.get_args(TableItem.model_fields["label"].annotation)
    assert TABLE_KINDS == {member.value for member in declared}
    assert TABLE_KINDS <= DOC_ITEM_LABELS

    with_data = {
        name
        for name, value in vars(document_module).items()
        if inspect.isclass(value) and hasattr(value, "model_fields") and "data" in value.model_fields
    }
    assert with_data == {"TableItem"}
    assert TableItem.model_fields["data"].annotation is TableData
    # And a ``TableItem`` always has one: the field is required, so the "table
    # label, no table data" condition is reachable only by unmaking a real object
    # (below) or through a provider stand-in.
    assert TableItem.model_fields["data"].is_required()


def test_powerpoint_speaker_notes_arrive_on_the_notes_layer_and_are_kept() -> None:
    payload = notes_pptx()

    # Docling's default iteration drops the note before anything can record it.
    plain = DocumentConverter(allowed_formats=[InputFormat.PPTX]).convert(
        _stream(payload, "deck.pptx"), raises_on_error=False
    )
    assert SPEAKER_NOTE not in [getattr(item, "text", None) for item, _ in plain.document.iterate_items()]

    result = DoclingDocumentParser().parse(payload, source_name="deck.pptx")

    notes = [one for one in result.document.elements if one.content_layer == "notes"]
    assert [one.text for one in notes] == [SPEAKER_NOTE]
    assert SPEAKER_NOTE in result.document.text
    assert result.call.content_layers_present == ("body", "notes")
    assert result.call.input_format == FORMAT_PPTX
    # PPTX, unlike DOCX, does carry page geometry — still parser evidence only.
    assert result.call.coordinate_grade == PARSER_PAGE_COORDINATES
    assert result.call.page_count == 1
    assert [one.text for one in result.document.elements if one.content_layer == "body"][0] == SLIDE_TITLE


def test_furniture_text_is_not_silently_dropped() -> None:
    document = DoclingDocument(name="probe")
    document.add_text(label=DocItemLabel.PAGE_HEADER, text="EPA-HQ", content_layer=ContentLayer.FURNITURE)
    document.add_text(label=DocItemLabel.TEXT, text="body text")

    # Docling's own default would return the body alone.
    assert [getattr(item, "text", None) for item, _ in document.iterate_items()] == ["body text"]

    result = _parse_document(document)

    assert [(one.content_layer, one.text) for one in result.document.elements] == [
        ("furniture", "EPA-HQ"),
        ("body", "body text"),
    ]


def test_background_and_invisible_layer_text_is_kept_and_stays_in_its_own_layer() -> None:
    """The two layers no fixture reached before, built with the real enum members.

    ``background`` and ``invisible`` are the layers a document uses for content it
    does not display; Docling's default iteration hides both, so nothing could even
    record that they were there.
    """
    document = DoclingDocument(name="probe")
    document.add_text(label=DocItemLabel.TEXT, text="background text", content_layer=ContentLayer.BACKGROUND)
    document.add_text(label=DocItemLabel.TEXT, text="invisible text", content_layer=ContentLayer.INVISIBLE)
    document.add_text(label=DocItemLabel.TEXT, text="body text")

    assert [getattr(item, "text", None) for item, _ in document.iterate_items()] == ["body text"]

    result = _parse_document(document)

    assert [(one.content_layer, one.text) for one in result.document.elements] == [
        ("background", "background text"),
        ("invisible", "invisible text"),
        ("body", "body text"),
    ]
    assert result.call.content_layers_present == ("background", "body", "invisible")
    # Preserved text on any layer is still only preserved text: the layer is what a
    # consumer excludes on, and it is recorded on every element.
    assert all(one.text_usable for one in result.document.elements)


def test_real_heading_context_never_crosses_a_content_layer() -> None:
    document = DoclingDocument(name="probe")
    document.add_text(label=DocItemLabel.TITLE, text="Body Title")
    document.add_text(label=DocItemLabel.SECTION_HEADER, text="Furniture Heading", content_layer=ContentLayer.FURNITURE)
    document.add_text(label=DocItemLabel.TEXT, text="Body paragraph.")
    document.add_text(label=DocItemLabel.TEXT, text="Footer line.", content_layer=ContentLayer.FURNITURE)
    document.add_text(label=DocItemLabel.TEXT, text="Background line.", content_layer=ContentLayer.BACKGROUND)

    result = _parse_document(document)

    paths = {one.text: (one.content_layer, one.heading_path) for one in result.document.elements}
    assert paths["Body paragraph."] == ("body", ("Body Title",))
    # The furniture heading describes furniture, and the body title does not
    # describe the page footer.
    assert paths["Footer line."] == ("furniture", ("Furniture Heading",))
    assert paths["Furniture Heading"] == ("furniture", ())
    # A layer with no heading of its own carries none, rather than borrowing one.
    assert paths["Background line."] == ("background", ())


# --- exact table geometry ---------------------------------------------------


def test_a_real_merged_cell_and_real_separator_text_survive_exactly() -> None:
    result = DoclingDocumentParser().parse(merged_table_docx(), source_name="merged.docx")

    table = result.document.tables[0]
    assert [(cell.row_start, cell.row_end, cell.column_start, cell.column_end, cell.text) for cell in table.cells] == [
        (0, 1, 0, 2, "Spanning header"),
        (1, 2, 0, 1, "Line\twith\ttabs"),
        (1, 2, 1, 2, "two\nlines"),
    ]
    # A merged cell writes once, at its top-left; the rest of its span is empty.
    assert result.document.text == "Spanning header\t\nLine\twith\ttabs\ttwo\nlines"
    # The flat rendering cannot be split back into these cells; the cells can.
    assert table.serialization_ambiguous is True
    assert result.call.table_cell_count == 3

    # The pinned release really does carry both descriptions of one cell shape,
    # which is why the reader checks them against each other.
    assert {"start_row_offset_idx", "end_row_offset_idx", "start_col_offset_idx", "end_col_offset_idx"} <= set(
        TableCell.model_fields
    )
    assert {"row_span", "col_span", "text", "column_header", "row_header"} <= set(TableCell.model_fields)


def _cell(row: int, column: int, text: str, *, row_span: int = 1, col_span: int = 1) -> TableCell:
    return TableCell(
        start_row_offset_idx=row,
        end_row_offset_idx=row + row_span,
        start_col_offset_idx=column,
        end_col_offset_idx=column + col_span,
        row_span=row_span,
        col_span=col_span,
        text=text,
    )


def _table_document(*, rows: int, columns: int, cells: list[TableCell]) -> DoclingDocument:
    document = DoclingDocument(name="probe")
    document.add_table(data=TableData(num_rows=rows, num_cols=columns, table_cells=cells))
    return document


def test_the_pinned_table_cell_really_declares_its_header_flags_as_booleans() -> None:
    # Why the reader refuses a non-bool instead of coercing with ``bool()``: the
    # pinned release declares these fields ``bool``, so anything else is malformed
    # provider output, not a value to interpret.
    assert TableCell.model_fields["column_header"].annotation is bool
    assert TableCell.model_fields["row_header"].annotation is bool
    assert {"num_rows", "num_cols", "table_cells"} <= set(TableData.model_fields)


@pytest.mark.parametrize(
    ("rows", "columns", "cells", "message"),
    [
        # A cell reaching past the declared grid: two descriptions of one table
        # that disagree. The grid used to widen silently to fit it.
        (1, 1, [_cell(1, 0, "30 mg/L")], "lies outside the declared row and column counts"),
        (1, 1, [_cell(0, 1, "30 mg/L")], "lies outside the declared row and column counts"),
        # Merged rectangles that overlap without sharing a top-left. One position
        # cannot hold two values, and only the shared anchor used to be caught.
        (
            1,
            2,
            [_cell(0, 0, "Spanning header", col_span=2), _cell(0, 1, "collides")],
            "share one grid position",
        ),
        (
            2,
            1,
            [_cell(0, 0, "Tall cell", row_span=2), _cell(1, 0, "collides")],
            "share one grid position",
        ),
        # A grid that is empty in one dimension only. ``0xN`` and ``Nx0`` are real,
        # constructible ``TableData`` values whose product is zero, so an area
        # bound alone accepts every one of them — including a column count no
        # machine could lay out — and records the result as a checked table.
        (0, 5, [], "empty in one dimension only"),
        (5, 0, [], "empty in one dimension only"),
        (0, 10**12, [], "empty in one dimension only"),
        (10**12, 0, [], "empty in one dimension only"),
    ],
)
def test_a_real_table_whose_declared_shape_contradicts_its_cells_fails_closed(
    rows: int, columns: int, cells: list[TableCell], message: str
) -> None:
    with pytest.raises(DoclingParseError, match=message) as failure:
        _parse_document(_table_document(rows=rows, columns=columns, cells=cells))

    assert failure.value.call is not None
    assert failure.value.call.failure_reason == "malformed_element"


@pytest.mark.parametrize(
    ("rows", "columns", "message", "reason"),
    [
        # One declared dimension past its own bound beside a small one: the product
        # bound alone lets this through whenever the other dimension is 0 or 1.
        (MAX_TABLE_DIMENSION + 1, 1, "past the recorded dimension bound", "table_dimension_limit"),
        (1, MAX_TABLE_DIMENSION + 1, "past the recorded dimension bound", "table_dimension_limit"),
        # And the declared area, checked before the no-cells shortcut so a grid
        # declared enormous and delivered empty cannot pass as a checked table.
        (MAX_TABLE_DIMENSION, 2, "recorded per-table cell bound", "table_cell_limit"),
    ],
)
def test_a_real_grid_declared_past_a_recorded_bound_names_the_bound_it_reached(
    rows: int, columns: int, message: str, reason: str
) -> None:
    # These are this adapter's own limits, not malformed provider output, so the
    # receipt names the bound rather than blaming the release.
    with pytest.raises(DoclingParseError, match=message) as failure:
        _parse_document(_table_document(rows=rows, columns=columns, cells=[]))

    call = failure.value.call
    assert call is not None
    assert call.failure_reason == reason
    assert (call.policy.max_table_dimension, call.policy.max_cells_per_table) == (
        MAX_TABLE_DIMENSION,
        MAX_CELLS_PER_TABLE,
    )


def test_the_pinned_table_cell_always_declares_both_spans_and_a_missing_one_fails_closed() -> None:
    """Requiring the spans cannot refuse valid provider output, and catches invalid.

    ``row_span`` and ``col_span`` carry defaults on the pinned ``TableCell``, so a
    real cell always has both — which is exactly why defaulting an absent one to
    ``1`` was wrong: it manufactured the agreement the offsets/span check exists to
    verify, and a merged cell that lost its span read as unmerged.
    """
    assert (TableCell.model_fields["row_span"].default, TableCell.model_fields["col_span"].default) == (1, 1)
    assert not TableCell.model_fields["row_span"].is_required()
    intact = _cell(0, 0, "Spanning header", col_span=2)
    assert (intact.row_span, intact.col_span) == (1, 2)

    for span in ("row_span", "col_span"):
        cell = _cell(0, 0, "Spanning header", col_span=2)
        delattr(cell, span)
        assert not hasattr(cell, span)

        with pytest.raises(DoclingParseError, match="table cell geometry is missing or invalid") as failure:
            _parse_document(_table_document(rows=1, columns=2, cells=[cell]))

        assert failure.value.call is not None
        assert failure.value.call.failure_reason == "malformed_element"


def test_a_real_table_labeled_item_carrying_no_table_data_fails_closed() -> None:
    """A ``TableItem`` has no ``.text`` at all, so unreadable ``TableData`` is total loss.

    ``TableItem.data`` is required, so this is reached by unmaking a real object —
    but the loss it guards against is the release blocker itself: mapped as an
    ordinary element, this document parses to an omission and no cell values, with
    nothing in the record saying a grid went missing.
    """
    document = DoclingDocument(name="probe")
    document.add_text(label=DocItemLabel.CAPTION, text="Table 1. Effluent limits")
    table = document.add_table(data=TableData(num_rows=1, num_cols=1, table_cells=[_cell(0, 0, "30 mg/L")]))
    delattr(table, "data")

    with pytest.raises(DoclingParseError, match="table-labeled item carries no table data") as failure:
        _parse_document(document)

    assert failure.value.call is not None
    assert failure.value.call.failure_reason == "malformed_element"


def test_a_real_item_carrying_table_data_under_a_non_table_label_fails_closed() -> None:
    """The other direction: a grid must not appear under a label that is not a table.

    ``TableItem.label`` is validated at construction but not on assignment, so a
    real, fully populated ``TableItem`` can end up carrying a ``picture`` label.
    Reading it as a table would publish a ``ParsedTable`` — grid, serialization,
    cell records — for an item the provider does not describe as one.
    """
    document = DoclingDocument(name="probe")
    document.add_text(label=DocItemLabel.CAPTION, text="Table 1. Effluent limits")
    table = document.add_table(data=TableData(num_rows=1, num_cols=1, table_cells=[_cell(0, 0, "30 mg/L")]))
    table.label = DocItemLabel.PICTURE
    assert isinstance(table, TableItem) and table.data.table_cells

    with pytest.raises(DoclingParseError, match="table data under a label that is not a table") as failure:
        _parse_document(document)

    assert failure.value.call is not None
    assert failure.value.call.failure_reason == "malformed_element"


def test_a_real_document_index_table_is_mapped_as_the_table_it_is() -> None:
    # The table-bearing label no fixture reached before, and the reason
    # ``TABLE_KINDS`` is a set rather than one label.
    document = DoclingDocument(name="probe")
    document.add_table(
        data=TableData(num_rows=1, num_cols=2, table_cells=[_cell(0, 0, "Part 401"), _cell(0, 1, "Scope")]),
        label=DocItemLabel.DOCUMENT_INDEX,
    )

    result = _parse_document(document)

    assert [one.kind for one in result.document.elements] == [DocItemLabel.DOCUMENT_INDEX.value]
    assert [cell.text for cell in result.document.tables[0].cells] == ["Part 401", "Scope"]
    assert result.document.text == "Part 401\tScope"


def test_a_real_caption_reference_that_cannot_describe_a_relationship_fails_closed() -> None:
    document = DoclingDocument(name="probe")
    caption = document.add_text(label=DocItemLabel.CAPTION, text="Table 1. Effluent limits")
    table = document.add_table(data=TableData(num_rows=1, num_cols=1, table_cells=[_cell(0, 0, "30 mg/L")]))

    # The valid link first, with the real provider objects.
    table.captions = [caption.get_ref()]
    assert _parse_document(document).document.tables[0].caption_refs == (caption.self_ref,)

    # Then a reference to an element the document does not contain. Recording it
    # would publish a link no consumer can follow.
    table.captions = [_reference("#/texts/99")]
    with pytest.raises(DoclingParseError, match="resolves to no emitted element") as dangling:
        _parse_document(document)
    assert dangling.value.call is not None
    assert dangling.value.call.failure_reason == "malformed_element"

    # An item that names itself: not a relationship between two elements at all.
    table.captions = [table.get_ref()]
    with pytest.raises(DoclingParseError, match="itself as its own caption"):
        _parse_document(document)

    # And one edge written twice in a single list, which is still one edge.
    table.captions = [caption.get_ref(), caption.get_ref()]
    with pytest.raises(DoclingParseError, match="names one caption more than once"):
        _parse_document(document)


def test_two_real_tables_may_share_one_caption_target() -> None:
    """``captions`` is a plain ``list[RefItem]``; nothing in the release makes it exclusive.

    Treating a caption as owned by one item refused a document the pinned
    docling-core produces without complaint, and the adapter has no standing to
    decide which of two tables a shared caption "really" belongs to.
    """
    document = DoclingDocument(name="probe")
    caption = document.add_text(label=DocItemLabel.CAPTION, text="Table 1. Effluent limits")
    shared = caption.get_ref()
    first = document.add_table(data=TableData(num_rows=1, num_cols=1, table_cells=[_cell(0, 0, "30 mg/L")]))
    second = document.add_table(data=TableData(num_rows=1, num_cols=1, table_cells=[_cell(0, 0, "45 mg/L")]))
    first.captions = [shared]
    second.captions = [shared]

    result = _parse_document(document)

    assert [table.caption_refs for table in result.document.tables] == [
        (caption.self_ref,),
        (caption.self_ref,),
    ]
    assert result.call.table_count == 2


def test_add_table_accepts_an_ordinary_text_item_as_a_caption_and_the_link_is_kept() -> None:
    """The pinned signature is ``caption: Optional[Union[TextItem, RefItem]]``.

    Any ``TextItem`` — whatever its label — and any bare reference are valid
    captions to the release, so restricting the target to a ``caption``-labelled
    element refused links ``add_table`` itself creates. What may be *quoted* as a
    caption is a ``source.py`` decision made against locked bytes; the adapter
    records the edge the provider declared and grades nothing.
    """
    signature = inspect.signature(DoclingDocument.add_table)
    assert "TextItem" in str(signature.parameters["caption"].annotation)

    document = DoclingDocument(name="probe")
    ordinary = document.add_text(label=DocItemLabel.TEXT, text="Ordinary body text.")
    footnote = document.add_text(label=DocItemLabel.FOOTNOTE, text="1. See § 1.2.")
    table = document.add_table(
        data=TableData(num_rows=1, num_cols=1, table_cells=[_cell(0, 0, "30 mg/L")]),
        caption=ordinary,
    )
    table.captions.append(footnote.get_ref())

    result = _parse_document(document)

    assert result.document.tables[0].caption_refs == (ordinary.self_ref, footnote.self_ref)
    kinds = {one.parser_ref: one.kind for one in result.document.elements}
    assert [kinds[ref] for ref in result.document.tables[0].caption_refs] == ["text", "footnote"]


def test_a_cell_whose_span_contradicts_its_offsets_fails_closed() -> None:
    document = DoclingDocument(name="probe")
    document.add_table(
        data=TableData(
            num_rows=1,
            num_cols=1,
            table_cells=[
                TableCell(
                    start_row_offset_idx=0,
                    end_row_offset_idx=1,
                    start_col_offset_idx=0,
                    end_col_offset_idx=1,
                    col_span=3,
                    text="30 mg/L",
                )
            ],
        )
    )

    with pytest.raises(DoclingParseError, match="spans disagree with its row and column offsets") as failure:
        _parse_document(document)

    assert failure.value.call is not None
    assert failure.value.call.failure_reason == "malformed_element"


# --- exact docling-core behavior the adapter depends on ---------------------


def test_a_formula_item_may_carry_its_source_only_in_orig() -> None:
    document = DoclingDocument(name="probe")
    empty_text = document.add_text(label=DocItemLabel.FORMULA, text="", orig=FORMULA_SOURCE)
    same_text = document.add_text(label=DocItemLabel.FORMULA, text="a+b")

    assert (empty_text.text, empty_text.orig) == ("", FORMULA_SOURCE)
    # When a caller gives no ``orig``, docling-core copies ``text`` into it.
    assert (same_text.text, same_text.orig) == ("a+b", "a+b")

    result = _parse_document(document)

    assert [one.text for one in result.document.elements] == [FORMULA_SOURCE, "a+b"]
    assert [one.content_source for one in result.document.elements] == [CONTENT_FROM_ORIG, CONTENT_FROM_TEXT]


def test_captions_stay_provider_elements_and_tables_keep_the_link() -> None:
    document = DoclingDocument(name="probe")
    caption = document.add_text(label=DocItemLabel.CAPTION, text="Table 1. Effluent limits")
    document.add_table(
        data=TableData(
            num_rows=1,
            num_cols=1,
            table_cells=[
                TableCell(
                    start_row_offset_idx=0,
                    end_row_offset_idx=1,
                    start_col_offset_idx=0,
                    end_col_offset_idx=1,
                    text="30 mg/L",
                )
            ],
        ),
        caption=caption,
    )

    result = _parse_document(document)

    assert [one.kind for one in result.document.elements] == ["caption", "table"]
    assert result.document.text == "Table 1. Effluent limits\n\n30 mg/L"
    assert result.document.tables[0].caption_refs == ("#/texts/0",)


# --- the loaded Office-only path -------------------------------------------


def test_the_loaded_converter_builds_the_office_only_model_free_path() -> None:
    built = DoclingDocumentParser()

    assert built.supported_formats == SUPPORTED_FORMATS
    options = _converter(built).format_to_options
    assert set(options) == {InputFormat.DOCX, InputFormat.PPTX, InputFormat.XLSX}
    assert isinstance(options[InputFormat.DOCX], WordFormatOption)
    assert isinstance(options[InputFormat.PPTX], PowerpointFormatOption)
    assert isinstance(options[InputFormat.XLSX], ExcelFormatOption)

    for input_format in options:
        applied = options[input_format].pipeline_options
        assert options[input_format].pipeline_cls is SimplePipeline
        assert type(applied) is ConvertPipelineOptions
        assert applied.enable_remote_services is False
        assert applied.allow_external_plugins is False
        assert applied.do_picture_classification is False
        assert applied.do_picture_description is False
        assert applied.do_chart_extraction is False
        # Nothing this path holds could load a model, so nothing needs a store.
        assert applied.artifacts_path is None
    # Each format holds its own options object: none can be reconfigured through
    # another's.
    assert len({id(options[one].pipeline_options) for one in options}) == 3


def test_simple_pipeline_really_ignores_the_document_timeout_the_policy_calls_unenforced() -> None:
    # The honesty this record depends on: ``document_timeout`` is a field on the
    # options SimplePipeline receives, and nothing in SimplePipeline reads it.
    assert "document_timeout" in ConvertPipelineOptions.model_fields
    assert ConvertPipelineOptions().document_timeout is None
    assert "document_timeout" not in inspect.getsource(SimplePipeline)

    policy = DoclingDocumentParser().policy
    assert (policy.document_timeout_enforced, policy.page_limit_enforced) == (False, False)


def test_the_policy_records_the_mapping_identity_and_the_bounds_it_really_enforces() -> None:
    """What the receipt claims about this parse, checked against the real converter.

    Three of these are honest negatives. ``SimplePipeline.execute`` runs the
    backend to completion in one call, so no wall-clock, CPU, memory, or
    archive-expansion bound can be applied from inside this adapter; the process
    containment that would hold them belongs to ``source.py``.
    """
    built = DoclingDocumentParser()

    assert built.policy.mapping_revision == ADAPTER_MAPPING_REVISION
    assert built.policy.table_serialization == TABLE_SERIALIZATION
    assert built.policy.text_offsets == PARSED_TEXT_OFFSETS
    # The mapping revision is nameable in the identity, not only hashed into it.
    assert built.parser_id == (
        f"docling:{DOCLING_VERSION}:docling-core:{DOCLING_CORE_VERSION}"
        f":{ADAPTER_MAPPING_REVISION}:{built.policy_digest[:16]}"
    )

    # Every bound the adapter really enforces is in the record it publishes, each
    # at the scope its name states, so every count in a receipt has a comparator.
    assert (built.policy.max_items, built.policy.max_tables) == (MAX_ITEMS, MAX_TABLES)
    assert (built.policy.max_cells_per_table, built.policy.max_mapped_characters) == (
        MAX_CELLS_PER_TABLE,
        MAX_MAPPED_CHARACTERS,
    )
    assert (built.policy.max_total_table_cells, built.policy.max_table_dimension) == (
        MAX_TOTAL_TABLE_CELLS,
        MAX_TABLE_DIMENSION,
    )
    assert (built.policy.max_heading_level, built.policy.max_reference_chars) == (
        MAX_HEADING_LEVEL,
        MAX_REFERENCE_CHARS,
    )
    assert (built.policy.max_caption_refs_per_item, built.policy.max_total_caption_refs) == (
        MAX_CAPTION_REFS_PER_ITEM,
        MAX_TOTAL_CAPTION_REFS,
    )
    assert (built.policy.max_regions_per_item, built.policy.max_total_regions) == (
        MAX_REGIONS_PER_ITEM,
        MAX_TOTAL_REGIONS,
    )
    assert built.policy.max_provider_errors == MAX_PROVIDER_ERRORS
    # Retained cell characters are their own bound at their own scope: mapped text
    # is what a consumer receives, and a whitespace-only table keeps every cell
    # while contributing none of it.
    assert built.policy.max_table_cell_characters == MAX_TABLE_CELL_CHARACTERS
    # And the ones it cannot: claimed nowhere, recorded as unenforced. No
    # wall-clock, CPU, memory, or archive-expansion bound is applied from inside
    # a library call, and ``source.py``'s process gate is where they belong.
    assert built.policy.process_containment_enforced is False
    source = inspect.getsource(SimplePipeline)
    assert "signal" not in source and "resource" not in source and "alarm" not in source


def test_the_heading_level_bound_is_the_level_the_pinned_release_declares() -> None:
    """``MAX_HEADING_LEVEL`` is docling-core's own ceiling, not a number we chose.

    It is also what bounds a recorded ``heading_path``: the per-layer heading stack
    keeps at most one entry per distinct level, so an unbounded level would let one
    document grow every element's heading path without bound.
    """
    metadata = SectionHeaderItem.model_fields["level"].metadata
    assert {type(entry).__name__: getattr(entry, "ge", None) or getattr(entry, "le", None) for entry in metadata} == {
        "Ge": 1,
        "Le": MAX_HEADING_LEVEL,
    }

    document = DoclingDocument(name="probe")
    document.add_heading(text="Deepest heading the release allows", level=MAX_HEADING_LEVEL)
    document.add_text(label=DocItemLabel.TEXT, text="Body under it.")

    result = _parse_document(document)

    assert result.document.elements[-1].heading_path == ("Deepest heading the release allows",)
    # The release refuses to build one deeper, which is what the adapter mirrors.
    with pytest.raises(ValueError):
        DoclingDocument(name="probe").add_heading(text="Too deep", level=MAX_HEADING_LEVEL + 1)


def test_every_reference_the_pinned_release_emits_fits_the_recorded_reference_bound() -> None:
    # The bound is this adapter's own — the pinned ``RefItem`` pattern constrains a
    # reference's shape but not its length — so it is held against references the
    # release really produces, with room to spare.
    result = DoclingDocumentParser().parse(notes_pptx(), source_name="deck.pptx")

    references = [one.parser_ref for one in result.document.elements]
    references += [one.parent_ref for one in result.document.elements if one.parent_ref]
    assert references
    assert max(len(reference) for reference in references) * 4 < MAX_REFERENCE_CHARS
    assert all(re.fullmatch(r"^#(?:/([\w-]+)(?:/(\d+))?)?$", reference) for reference in references)


def test_a_pdf_is_refused_by_name_before_the_provider_runs() -> None:
    recorded = RecordedConverter(office_converter())
    built = DoclingDocumentParser(converter=recorded)

    with pytest.raises(DoclingParseError, match="recognized but not implemented") as failure:
        built.parse(b"%PDF-1.7\n", source_name="rule.pdf")

    call = failure.value.call
    assert call is not None
    assert (call.failure_reason, call.provider_invoked, call.input_format) == ("format_not_implemented", False, "pdf")
    assert recorded.calls == []


def test_the_adapter_hands_the_provider_its_own_file_name_and_no_borrowed_limits() -> None:
    recorded = RecordedConverter(office_converter())
    built = DoclingDocumentParser(converter=recorded, max_source_bytes=10_000_000)

    result = built.parse(table_only_docx(), source_name="https://example.gov/private/limits.docx?token=abc")

    assert (result.call.source_name, result.call.source_name_sanitized) == ("limits.docx", True)
    assert result.call.provider_input_format == FORMAT_DOCX
    assert recorded.calls == [{"name": "source.docx", "options": {"raises_on_error": False}}]


def test_a_real_format_mismatch_fails_closed() -> None:
    # DOCX bytes announced as a spreadsheet. Docling sniffs the content, picks the
    # Word backend, and reports success — so the recorded format would have
    # described a pipeline that never ran. The adapter refuses instead.
    with pytest.raises(DoclingParseError, match="a different input format than the one detected") as failure:
        DoclingDocumentParser().parse(table_only_docx(), source_name="limits.xlsx")

    call = failure.value.call
    assert call is not None
    assert (call.conversion_status, call.failure_reason) == ("success", "format_mismatch")
    assert (call.input_format, call.provider_input_format) == ("xlsx", FORMAT_DOCX)


def test_an_over_limit_input_never_reaches_the_provider() -> None:
    payload = table_only_docx()

    with pytest.raises(DoclingParseError, match="exceed the recorded input limit") as failure:
        DoclingDocumentParser(max_source_bytes=len(payload) - 1).parse(payload, source_name="limits.docx")

    assert failure.value.call is not None
    assert failure.value.call.provider_invoked is False


# --- the provider scalars a record persists ---------------------------------


def _provenance(page_no: Any = 1, charspan: Any = (0, 5)) -> ProvenanceItem:
    """One real ``ProvenanceItem``, built exactly the way the pinned model declares it."""
    return ProvenanceItem(
        page_no=page_no,
        bbox=BoundingBox(l=0.0, t=1.0, r=1.0, b=0.0, coord_origin=CoordOrigin.BOTTOMLEFT),
        charspan=charspan,
    )


def _located_document(prov: ProvenanceItem) -> DoclingDocument:
    document = DoclingDocument(name="probe")
    document.add_text(label=DocItemLabel.TEXT, text="Located text.", prov=prov)
    return document


def test_the_pinned_provenance_item_bounds_neither_its_page_number_nor_its_character_span() -> None:
    """Why these two scalars need a bound of this adapter's own: the release has none.

    ``ProvenanceItem`` declares ``page_no: int`` and ``charspan: Tuple[int, int]``
    and validates neither for magnitude, so the pinned model happily accepts — and
    hands over — an integer of arbitrary precision. Left unbounded, that value
    landed verbatim on a frozen :class:`PageRegion` and went on to a receipt and a
    Parquet row, where no reader has a comparator for it and no fixed-width column
    can hold it.
    """
    enormous = 10**100
    unbounded = _provenance(page_no=enormous, charspan=(0, enormous))

    # The release itself objects to none of it, and keeps the exact integer.
    assert unbounded.page_no == enormous
    assert unbounded.charspan == (0, enormous)
    assert {name: field.annotation for name, field in ProvenanceItem.model_fields.items()}["page_no"] is int

    for prov, phrase, reason in (
        (_provenance(page_no=enormous), "page-number bound", "page_number_limit"),
        (_provenance(charspan=(0, enormous)), "character-index bound", "character_span_limit"),
    ):
        with pytest.raises(DoclingParseError, match=phrase) as failure:
            _parse_document(_located_document(prov))
        call = failure.value.call
        assert call is not None
        assert call.failure_reason == reason
        assert str(enormous) not in str(call.as_json_dict())


@pytest.mark.parametrize(
    ("prov", "expected"),
    [
        (_provenance(page_no=MAX_PAGE_NUMBER), ("page_number", MAX_PAGE_NUMBER)),
        (_provenance(charspan=(0, MAX_PROVENANCE_CHARACTER_INDEX)), ("char_end", MAX_PROVENANCE_CHARACTER_INDEX)),
    ],
    ids=["page number", "provenance character index"],
)
def test_a_real_provenance_scalar_exactly_at_its_bound_is_kept_verbatim(
    prov: ProvenanceItem, expected: tuple[str, int]
) -> None:
    # A ceiling, not a threshold, against the real model — and kept verbatim, since
    # provenance is parser evidence this adapter records rather than interprets.
    result = _parse_document(_located_document(prov))

    field, value = expected
    assert getattr(result.document.elements[0].regions[0], field) == value


@pytest.mark.parametrize(
    ("prov", "phrase", "reason"),
    [
        (_provenance(page_no=MAX_PAGE_NUMBER + 1), "page-number bound", "page_number_limit"),
        (
            _provenance(charspan=(0, MAX_PROVENANCE_CHARACTER_INDEX + 1)),
            "character-index bound",
            "character_span_limit",
        ),
    ],
    ids=["page number", "provenance character index"],
)
def test_a_real_provenance_scalar_one_past_its_bound_names_the_bound_it_reached(
    prov: ProvenanceItem, phrase: str, reason: str
) -> None:
    with pytest.raises(DoclingParseError, match=phrase) as failure:
        _parse_document(_located_document(prov))

    call = failure.value.call
    assert call is not None
    # This adapter's own limit, so the receipt names it rather than blaming the release.
    assert call.failure_reason == reason
    assert (call.policy.max_page_number, call.policy.max_provenance_char_index) == (
        MAX_PAGE_NUMBER,
        MAX_PROVENANCE_CHARACTER_INDEX,
    )
    # And still not an operational page limit: the conversion is already over.
    assert call.policy.page_limit_enforced is False


def _nested_document(depth: int) -> DoclingDocument:
    """A real document whose one text item sits ``depth`` groups below the body."""
    document = DoclingDocument(name="probe")
    parent = None
    for index in range(depth):
        parent = document.add_group(name=f"group-{index}", parent=parent)
    document.add_text(label=DocItemLabel.TEXT, text="Deeply nested text.", parent=parent)
    return document


def test_a_real_tree_level_is_the_depth_the_release_reports_and_carries_a_stated_ceiling() -> None:
    """``iterate_items`` yields ``(item, level)`` and the level is persisted verbatim.

    It is a provider-controlled integer like any other, so it is held to a stated
    ceiling — and the real depths an Office document produces sit orders of
    magnitude under it, which is what makes the bound a ceiling rather than a
    limit on real output.
    """
    document = _nested_document(3)
    levels = [level for _, level in document.iterate_items()]

    result = _parse_document(document)

    assert levels == [4]
    assert [one.tree_level for one in result.document.elements] == levels
    assert result.call.policy.max_tree_depth == MAX_TREE_DEPTH

    # Real Office output is nowhere near the bound.
    deck = DoclingDocumentParser().parse(notes_pptx(), source_name="deck.pptx")
    assert max(one.tree_level for one in deck.document.elements) * 100 < MAX_TREE_DEPTH


def test_a_real_tree_level_past_the_recorded_bound_names_the_bound_it_reached(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    # Driven with a lowered bound over a real document, because the pinned
    # release's own recursive iteration cannot reach the real ceiling — see below.
    monkeypatch.setattr(adapter_module, "MAX_TREE_DEPTH", 2)

    with pytest.raises(DoclingParseError, match="deeper than the recorded tree-depth bound") as failure:
        _parse_document(_nested_document(3))

    call = failure.value.call
    assert call is not None
    assert (call.failure_reason, call.policy.max_tree_depth) == ("tree_depth_limit", 2)


def test_the_releases_own_iteration_gives_out_before_the_tree_depth_bound_can_be_reached() -> None:
    """An honest limit of the real ceiling, recorded rather than papered over.

    ``iterate_items`` recurses, so a document nested past roughly a thousand groups
    exhausts Python's stack inside the pinned release — before this adapter's
    ``MAX_TREE_DEPTH`` could refuse it. That is not an escape hatch: the failure is
    caught like any other provider failure and recorded with a complete, bounded
    receipt, rather than propagating as a bare ``RecursionError``.
    """
    unreachable = _nested_document(MAX_TREE_DEPTH + 1)
    with pytest.raises(RecursionError):
        list(unreachable.iterate_items())

    with pytest.raises(DoclingParseError, match="docling document could not be read") as failure:
        _parse_document(unreachable)

    call = failure.value.call
    assert call is not None
    assert (call.status, call.failure_reason) == ("failed", "provider_error")
    assert call.error_type == "RecursionError"


# --- provider exception type identity, against the real releases -------------


def _sampled_provider_exception_names() -> tuple[set[str], set[str], set[str]]:
    """A *sample* of the exception classes the two pinned distributions define.

    Sampled, not exhaustive, and the name of this helper says so. Two deliberate
    narrowings: only modules whose dotted name contains ``exception`` or ``error``
    are opened, so an exception class defined beside ordinary code elsewhere in
    the tree is never seen; and a module that fails to import is passed over,
    because walking a package imports it and some pinned modules need optional
    extras this suite does not install.

    Returns the names found, the modules that were considered, and the modules
    that were skipped — the last so a caller can assert on what the sample missed
    rather than have it disappear silently. Traversing a package imports it and
    some pinned modules warn on import; that is noise about the release, not about
    this adapter, so it is silenced.
    """
    names: set[str] = set()
    considered: set[str] = set()
    skipped: set[str] = set()
    with warnings.catch_warnings():
        warnings.simplefilter("ignore")
        for package in (docling, docling_core):
            for module in pkgutil.walk_packages(package.__path__, f"{package.__name__}.", onerror=skipped.add):
                if "exception" not in module.name and "error" not in module.name:
                    continue
                considered.add(module.name)
                try:
                    loaded = importlib.import_module(module.name)
                except Exception:
                    skipped.add(module.name)
                    continue
                names |= {
                    value.__name__
                    for value in vars(loaded).values()
                    if inspect.isclass(value) and issubclass(value, BaseException)
                }
    return (names, considered, skipped)


def test_the_sampled_provider_exception_names_survive_the_recorded_bound() -> None:
    """Bounding a class name must not destroy the real identities it exists to keep.

    The bound and the secret heuristic are only worth having if an ordinary
    provider failure still records *which* failure it was.

    The provider half of this is a **sample**, not an inventory: only modules
    named for exceptions or errors are opened, and one that will not import is
    skipped. So this asserts on the anchors it must find, reports what the sample
    passed over, and claims nothing about exception classes it never looked at.
    The unconditional guarantee is elsewhere — :func:`bounded_error_type` replaces
    any name it cannot keep, whatever the class, so an unsampled one is bounded
    too; it just may be recorded as the fallback rather than by name.
    """
    provider, considered, skipped = _sampled_provider_exception_names()

    # The sample really opened something, and really found the anchors in it.
    assert considered - skipped
    assert {"ConversionError", "OperationNotAllowed"} <= provider
    assert not {name for name in skipped if "docling.exceptions" == name}, sorted(skipped)
    # What the sample passed over is visible rather than silent; nothing below is
    # a claim about these.
    if skipped:
        print(f"unsampled provider modules: {sorted(skipped)}")  # noqa: T201

    standard = {
        "TypeError",
        "ValueError",
        "RuntimeError",
        "AttributeError",
        "KeyError",
        "IndexError",
        "OSError",
        "MemoryError",
        "RecursionError",
        "UnicodeEncodeError",
        "UnicodeDecodeError",
    }
    for name in sorted(provider | standard):
        assert len(name) <= MAX_ERROR_TYPE_CHARS, name
        assert bounded_error_type(name) == name, name


class _RaisingConverter:
    """A converter that fails the way a provider does: by raising its own class."""

    def __init__(self, error: BaseException) -> None:
        self._error = error
        self.calls: list[str] = []

    def convert(self, source: Any, **options: Any) -> Any:
        del options
        self.calls.append(Path(str(source)).name)
        raise self._error


def test_a_real_provider_exception_class_name_is_bounded_before_it_is_recorded() -> None:
    """A subclass of the release's own ``ConversionError``, named hostilely.

    Nothing stops a provider — or a plugin inside one — from raising a class built
    at runtime, and its ``__name__`` reached both the public failure message and a
    persisted ``error_type``.
    """
    hostile = type("SecretAkiaX" + "X" * 100_000, (ConversionError,), {})
    converter = _RaisingConverter(hostile("boom /secret/scan.docx"))

    with pytest.raises(DoclingParseError) as failure:
        DoclingDocumentParser(converter=converter).parse(table_only_docx(), source_name="limits.docx")

    call = failure.value.call
    assert call is not None
    assert (call.failure_reason, call.error_type) == ("provider_error", FALLBACK_ERROR_TYPE)
    assert str(failure.value) == f"docling conversion failed with {FALLBACK_ERROR_TYPE}"
    payload = str(call.as_json_dict())
    assert len(payload) < 5_000
    for fragment in ("XXXXXXXXXX", "Secret", "Akia", "scan.docx", "boom"):
        assert fragment not in str(failure.value)
        assert fragment not in payload

    # An ordinary provider failure still records which one it was.
    plain = _RaisingConverter(ConversionError("boom /secret/scan.docx"))
    with pytest.raises(DoclingParseError, match="conversion failed with ConversionError") as ordinary:
        DoclingDocumentParser(converter=plain).parse(table_only_docx(), source_name="limits.docx")
    assert ordinary.value.call is not None
    assert ordinary.value.call.error_type == "ConversionError"


# --- the mapped-character budget over real tables ---------------------------


def _watch_serialization(monkeypatch: pytest.MonkeyPatch) -> list[Any]:
    calls: list[Any] = []
    real = adapter_module._serialized_table_text

    def watched(reading: Any) -> str:
        calls.append(reading)
        return real(reading)

    monkeypatch.setattr(adapter_module, "_serialized_table_text", watched)
    return calls


def test_a_real_tables_serialized_length_is_computed_exactly_before_the_grid_exists() -> None:
    """The refusal is only as good as the arithmetic it is made on.

    Held against a real merged cell and real cell text carrying the
    serialization's own separators, recomputed from the persisted
    :class:`ParsedTable` — which holds exactly the values the length was computed
    from before any grid was allocated.
    """
    result = DoclingDocumentParser().parse(merged_table_docx(), source_name="merged.docx")

    table = result.document.tables[0]
    element = result.document.elements[table.element_ordinal]
    reading = adapter_module._TableReading(
        row_count=table.row_count,
        column_count=table.column_count,
        cells=table.cells,
        caption_refs=table.caption_refs,
    )

    assert adapter_module._serialized_table_length(reading) == len(element.text)
    assert element.text == adapter_module._serialized_table_text(reading)


def test_a_real_table_over_the_remaining_budget_never_reaches_its_serialization(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    """The allocation is the harm, so the refusal happens ahead of it.

    A real 2x2 ``TableData`` whose serialization is one character past the bound:
    the grid is never built and the joined document is never assembled.
    """
    cells = [_cell(row, column, text) for row, values in enumerate(TABLE_ROWS) for column, text in enumerate(values)]
    document = _table_document(rows=2, columns=2, cells=cells)

    monkeypatch.setattr(adapter_module, "MAX_MAPPED_CHARACTERS", len(TABLE_TEXT))
    at_bound = _watch_serialization(monkeypatch)
    assert _parse_document(document).document.text == TABLE_TEXT
    assert len(at_bound) == 1

    monkeypatch.setattr(adapter_module, "MAX_MAPPED_CHARACTERS", len(TABLE_TEXT) - 1)
    serialized = _watch_serialization(monkeypatch)
    with pytest.raises(DoclingParseError, match="table text exceeds the recorded mapped-character bound") as failure:
        _parse_document(document)

    assert failure.value.call is not None
    assert failure.value.call.failure_reason == "character_limit"
    assert serialized == []


def test_a_real_table_after_earlier_text_is_measured_against_what_is_left(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    # The budget a table gets is the document's remaining room, separator included
    # — not the whole bound, which would let legal-on-their-own tables past it
    # together.
    lead = "Effluent limits follow."
    document = DoclingDocument(name="probe")
    document.add_text(label=DocItemLabel.TEXT, text=lead)
    document.add_table(
        data=TableData(
            num_rows=2,
            num_cols=2,
            table_cells=[
                _cell(row, column, text) for row, values in enumerate(TABLE_ROWS) for column, text in enumerate(values)
            ],
        )
    )
    exact = len(lead) + len(ELEMENT_SEPARATOR) + len(TABLE_TEXT)

    monkeypatch.setattr(adapter_module, "MAX_MAPPED_CHARACTERS", exact)
    assert _parse_document(document).call.character_count == exact

    monkeypatch.setattr(adapter_module, "MAX_MAPPED_CHARACTERS", exact - 1)
    serialized = _watch_serialization(monkeypatch)
    with pytest.raises(DoclingParseError, match="table text exceeds") as failure:
        _parse_document(document)

    assert failure.value.call is not None
    assert failure.value.call.failure_reason == "character_limit"
    assert serialized == []


# --- an identifier this adapter cannot encode -------------------------------


@pytest.mark.parametrize("raw", ["\ud800", "\udfff", "\ud800rule.docx", "rule\udc00.docx"])
def test_a_real_surrogate_source_name_never_reaches_the_real_converter(
    monkeypatch: pytest.MonkeyPatch, raw: str
) -> None:
    """A ``str`` holding a lone surrogate has no UTF-8 encoding at all.

    ``source_name_sha256`` covers the caller's exact identifier in those bytes, so
    there is no honest identity to build a record from. Refused before it is
    hashed, before it is sanitized, and before the real converter — or the real
    filesystem — ever sees the call, as a stable :class:`DoclingParseError` rather
    than a raw ``UnicodeEncodeError`` out of the identity build.
    """
    recorded = RecordedConverter(office_converter())
    sanitizing: list[str] = []
    real_sanitizer = adapter_module.sanitized_source_name
    monkeypatch.setattr(
        adapter_module,
        "sanitized_source_name",
        lambda name: (sanitizing.append(name), real_sanitizer(name))[1],
    )

    with pytest.raises(DoclingParseError, match="not encodable as utf-8") as failure:
        DoclingDocumentParser(converter=recorded).parse(table_only_docx(), source_name=raw)

    assert failure.value.call is None
    assert isinstance(failure.value.__cause__, UnicodeEncodeError)
    assert sanitizing == []
    assert recorded.calls == []


def test_a_real_non_ascii_source_name_parses_and_is_hashed_as_utf8() -> None:
    # Refusing an unencodable name must not refuse an ordinary one: text outside
    # ASCII is text, and its digest covers the caller's exact value.
    raw = "règle-Ω-2026.docx"
    recorded = RecordedConverter(office_converter())

    result = DoclingDocumentParser(converter=recorded).parse(table_only_docx(), source_name=raw)

    assert result.call.source_name_sha256 == hashlib.sha256(raw.encode(SOURCE_NAME_ENCODING)).hexdigest()
    assert (result.call.source_name, result.call.source_name_sanitized) == ("r_gle-_-2026.docx", True)
    assert result.document.text == TABLE_TEXT
    # The adapter names the file it writes; an untrusted identifier never reaches
    # the filesystem, whatever alphabet it is in.
    assert [one["name"] for one in recorded.calls] == ["source.docx"]


def test_a_real_whitespace_only_table_costs_no_mapped_room_and_never_serializes(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    """A real ``TableData`` whose every cell is whitespace maps nothing at all.

    Its serialization would be a truthy run of tabs and newlines that ``_assemble``
    then drops, recording the blank-table omission — so reserving mapped-character
    room for it refused documents that really fit. The cells are still preserved,
    because they are the table's only values.
    """
    blank = [
        _cell(row, column, text)
        for row, values in enumerate(((" ", "\t"), ("\n", "  ")))
        for column, text in enumerate(values)
    ]
    document = DoclingDocument(name="probe")
    document.add_text(label=DocItemLabel.TEXT, text="Effluent limits follow.")
    document.add_table(data=TableData(num_rows=2, num_cols=2, table_cells=blank))

    monkeypatch.setattr(adapter_module, "MAX_MAPPED_CHARACTERS", len("Effluent limits follow."))
    serialized = _watch_serialization(monkeypatch)

    result = _parse_document(document)

    assert result.document.text == "Effluent limits follow."
    assert serialized == []
    assert [cell.text for cell in result.document.tables[0].cells] == [" ", "\t", "\n", "  "]
    assert [one.reason for one in result.document.omissions] == ["blank-table-cells"]
    assert result.document.elements[result.document.tables[0].element_ordinal].content_source == NO_CONTENT


def test_the_cells_a_real_blank_table_retains_are_bounded_by_the_limit_that_names_them(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    # Not charging omitted whitespace to the mapped budget must not leave the
    # retained values unbounded: they sit under the bound that says what it limits.
    blank = [_cell(0, column, text) for column, text in enumerate(("  ", "   "))]
    document = DoclingDocument(name="probe")
    document.add_text(label=DocItemLabel.TEXT, text="Effluent limits follow.")
    document.add_table(data=TableData(num_rows=1, num_cols=2, table_cells=blank))

    monkeypatch.setattr(adapter_module, "MAX_TABLE_CELL_CHARACTERS", 5)
    assert _parse_document(document).call.status == "completed"

    monkeypatch.setattr(adapter_module, "MAX_TABLE_CELL_CHARACTERS", 4)
    with pytest.raises(DoclingParseError, match="recorded table-cell character bound") as failure:
        _parse_document(document)

    assert failure.value.call is not None
    assert failure.value.call.failure_reason == "table_cell_character_limit"


def test_every_scalar_the_pinned_releases_hand_over_is_an_exact_builtin_type() -> None:
    """Why demanding exact built-in scalars costs the pinned releases nothing.

    The adapter refuses an ``int`` or ``str`` *subtype* at every provider boundary,
    because a subtype answers the comparisons, hashes, and lengths every bound is
    made of. That is only affordable because the pinned releases emit plain
    built-ins everywhere the adapter reads — pydantic validates to ``int``,
    ``float``, ``str``, and ``bool``, and a ``str`` enum's ``.value`` is an
    ordinary string. This walks the whole read surface of three real Office
    documents and holds every value to its exact type, so a release that started
    handing back a subtype would be caught here rather than at a refused parse.
    """
    seen: set[str] = set()

    def record(name: str, value: Any, expected: type) -> None:
        assert type(value) is expected, f"{name} is a {type(value)!r}, not exactly {expected!r}"
        seen.add(name)

    documents = (
        (table_only_docx(), "t.docx"),
        (headings_docx(), "h.docx"),
        (notes_pptx(), "n.pptx"),
        (limits_xlsx(), "l.xlsx"),
    )
    for payload, name in documents:
        conversion = office_converter().convert(_stream(payload, name), raises_on_error=False)
        record("status", conversion.status.value, str)
        record("input format", conversion.input.format.value, str)
        record("page count", len(conversion.document.pages), int)
        for item, level in conversion.document.iterate_items(included_content_layers=set(ContentLayer)):
            record("tree level", level, int)
            record("label", item.label.value, str)
            record("content layer", item.content_layer.value, str)
            record("self reference", item.self_ref, str)
            if getattr(item, "parent", None) is not None:
                record("parent reference", item.parent.cref, str)
            if isinstance(getattr(item, "text", None), str):
                record("item text", item.text, str)
            if getattr(item, "level", None) is not None:
                record("heading level", item.level, int)
            for provenance in getattr(item, "prov", ()) or ():
                record("page number", provenance.page_no, int)
                record("charspan start", provenance.charspan[0], int)
                record("charspan end", provenance.charspan[1], int)
                record("coordinate origin", provenance.bbox.coord_origin.value, str)
                for edge in ("l", "t", "r", "b"):
                    record("bounding box edge", getattr(provenance.bbox, edge), float)
            data = getattr(item, "data", None)
            for cell in getattr(data, "table_cells", ()) or ():
                record("cell text", cell.text, str)
                record("cell row offset", cell.start_row_offset_idx, int)
                record("cell column offset", cell.start_col_offset_idx, int)
                record("cell row span", cell.row_span, int)
                record("cell column span", cell.col_span, int)
                record("cell header flag", cell.column_header, bool)
            if data is not None and hasattr(data, "table_cells"):
                record("declared row count", data.num_rows, int)
                record("declared column count", data.num_cols, int)

    # And the sweep really covered the whole surface rather than the easy half.
    assert seen == {
        "status", "input format", "page count", "tree level", "label", "content layer", "self reference",
        "parent reference", "item text", "heading level", "page number", "charspan start", "charspan end",
        "coordinate origin", "bounding box edge", "cell text", "cell row offset", "cell column offset",
        "cell row span", "cell column span", "cell header flag", "declared row count", "declared column count",
    }  # fmt: skip


def test_a_real_provider_property_raising_the_public_error_lends_it_no_message() -> None:
    """``DoclingParseError`` is public and importable, so provider code may raise it.

    A pydantic property, a lazily built model, or an iterator can raise this
    module's own error class carrying a message that quotes the source document.
    Held here against a real ``DoclingDocument``, whose iteration the adapter
    drives directly.
    """
    secret = "AKIAIOSFODNN7EXAMPLE /var/secrets/scan.docx " + "s" * 50_000

    class RaisingDocument:
        pages: dict[int, Any] = {}

        def iterate_items(self, **options: Any) -> Any:
            del options
            raise DoclingParseError(secret)

    built = DoclingDocumentParser(converter=_FixedConversion(typing.cast(Any, RaisingDocument())))

    with pytest.raises(DoclingParseError) as failure:
        built.parse(b"real docling-core document", source_name="probe.docx")

    assert str(failure.value) == "docling document could not be read: DoclingParseError"
    assert failure.value.call is not None
    assert failure.value.call.failure_reason == "provider_error"
    for fragment in ("AKIA", "secrets", "scan.docx", "ssssssssss"):
        assert fragment not in str(failure.value)


# --- helpers ---------------------------------------------------------------


def _reference(cref: str) -> RefItem:
    """Build a bare ``RefItem`` the way the pinned model declares it.

    ``RefItem.cref`` carries the alias ``$ref``, which is not a Python identifier,
    so the model is populated through its alias here rather than by keyword — the
    same value either way, and the only spelling a type checker can follow.
    """
    return RefItem.model_validate({"$ref": cref})


def _stream(payload: bytes, name: str) -> Any:
    from docling_core.types.io import DocumentStream  # ty: ignore[unresolved-import]

    return DocumentStream(name=name, stream=io.BytesIO(payload))


def _converter(built: DoclingDocumentParser) -> Any:
    """Reach the one private provider object, which the adapter's interface never names.

    Through ``__dict__``, because that is what reaching a private attribute in
    Python is: the boundary the adapter holds is that its protocol, its public
    attributes, and the records it returns name no provider type — not that an
    object inside it is unreachable.
    """
    return built.__dict__["_converter"]


class _FixedConversion:
    """A conversion result over a document built with the real docling-core types.

    Only the four fields the adapter reads, and every provider value in them —
    the status enum, the input format enum, and the ``DoclingDocument`` itself —
    is the real one.
    """

    def __init__(self, document: DoclingDocument) -> None:
        self.status = ConversionStatus.SUCCESS
        self.errors: list[Any] = []
        self.document = document
        self.input = SimpleNamespace(format=InputFormat.DOCX)

    def convert(self, source: Any, **options: Any) -> Any:
        del source, options
        return self


def _parse_document(document: DoclingDocument) -> Any:
    """Run the adapter over a document built directly with real docling-core types."""
    built = DoclingDocumentParser(converter=_FixedConversion(document))
    return built.parse(b"real docling-core document", source_name="probe.docx")
